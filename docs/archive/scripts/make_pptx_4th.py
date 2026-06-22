"""
렉처노트 4차 발표 PPTX 생성 스크립트
design.md 디자인 시스템 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_DIR = r"g:\내 드라이브\02_Study\01_Univ\2026학년도 1학기\06_경영전략\조별\lecture_note_project"
IMG_DIR = os.path.join(BASE_DIR, "_outputs", "4차 발표")
OUTPUT_PATH = os.path.join(BASE_DIR, "_history", "렉쳐노트_4차발표.pptx")

# ─── 컬러 팔레트 ────────────────────────────────────────────
NAVY        = RGBColor(0x1A, 0x2D, 0x5E)
NAVY_CARD   = RGBColor(0x22, 0x38, 0x72)  # 카드용 약간 밝은 Navy
SKY_BLUE    = RGBColor(0x4B, 0x82, 0xC5)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
GREEN_BG    = RGBColor(0xE8, 0xF8, 0xEE)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_BG   = RGBColor(0xFE, 0xF3, 0xC7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF2, 0xF7)
CHARCOAL    = RGBColor(0x1C, 0x1C, 0x1E)
GRAY        = RGBColor(0x6E, 0x6E, 0x73)
RED         = RGBColor(0xE7, 0x4C, 0x3C)

# 슬라이드 크기 (16:9)
SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# ─── 헬퍼 함수 ──────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, line_pt=0.75):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt2(line_pt)
    else:
        shape.line.fill.background()
    return shape

def oval(slide, x, y, w, h, fill=None):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def tb(slide, x, y, w, h, text, size=16, color=WHITE, bold=False, italic=False,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "맑은 고딕"
    return box

def pic(slide, img_path, x, y, w=None, h=None):
    if not os.path.exists(img_path):
        print(f"  [경고] 이미지 없음: {img_path}")
        return None
    if w and h:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
    elif w:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), height=Inches(h))

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def footer_skyblue(slide, text="경영전략 팀 프로젝트 | 9팀 4차 발표 | 2026.05"):
    rect(slide, 0, 6.92, 13.333, 0.58, fill=SKY_BLUE)
    tb(slide, 0, 6.93, 13.333, 0.56, text, size=13, color=WHITE, align=PP_ALIGN.CENTER)

def footer_navy(slide, text):
    rect(slide, 0, 6.92, 13.333, 0.58, fill=NAVY_CARD)
    tb(slide, 0.3, 6.93, 12.7, 0.56, text, size=13, color=WHITE, align=PP_ALIGN.CENTER)

def navy_bg(slide):
    rect(slide, 0, 0, 13.333, 7.5, fill=NAVY)

def slide_title_white(slide, title, subtitle=None):
    tb(slide, 0.5, 0.25, 12.3, 0.75, title, size=32, color=NAVY, bold=True)
    if subtitle:
        tb(slide, 0.5, 0.95, 12.3, 0.45, subtitle, size=17, color=GRAY)

def slide_title_navy(slide, title, subtitle=None):
    tb(slide, 0.5, 0.25, 12.3, 0.75, title, size=32, color=WHITE, bold=True)
    if subtitle:
        tb(slide, 0.5, 0.95, 12.3, 0.45, subtitle, size=17, color=SKY_BLUE)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)

# 좌측 세로선
rect(s, 0.85, 1.4, 0.07, 4.2, fill=SKY_BLUE)

# 차수 정보 (소형)
tb(s, 1.05, 1.25, 10, 0.45, "경영전략 9팀  |  4차 발표", size=16, color=GRAY)

# 영문 서비스명
tb(s, 1.05, 1.75, 11, 1.3, "LectureNote", size=78, color=WHITE, bold=True)

# 한국어 서비스명
tb(s, 1.05, 2.95, 10, 0.65, "렉처노트", size=38, color=SKY_BLUE, bold=True)

# 구분선
rect(s, 1.05, 3.68, 5.5, 0.06, fill=SKY_BLUE)

# 부제
tb(s, 1.05, 3.82, 11.5, 0.55,
   "알파 빌드 달성 — 웹 UI + 원클릭 녹음 완성",
   size=21, color=WHITE)
tb(s, 1.05, 4.42, 11.5, 0.45,
   "3차 발표에서 예고한 8주차 목표 완료",
   size=16, color=GRAY, italic=True)

# 팀 정보
tb(s, 1.05, 5.55, 11, 0.45,
   "이병헌 · 김도환 · 최철원  |  2026년 5월",
   size=16, color=GRAY)

footer_skyblue(s)

notes(s, """안녕하세요, 경영전략 9팀입니다.
저희는 강의 녹음·수업 자료·개인 필기를 결합해 AI가 통합 요약 노트를 자동 생성하는 서비스 '렉처노트'를 개발하고 있습니다.
이번 4차 발표에서는 3차 발표에서 약속드린 8주차 목표, 즉 사용자가 코드 없이 사용할 수 있는 웹 프론트엔드와 원클릭 녹음 기능이 완성된 알파 빌드를 직접 보여드리겠습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "발표 구성")

toc = [
    ("01", "SDLC 진행 현황",     "3차 대비 구현 단계 진척도"),
    ("02", "프로젝트 로드맵",     "8주차 알파 빌드 목표 달성"),
    ("03", "알파 빌드 주요 기능", "3가지 핵심 기능 완성"),
    ("04", "데모 — 강의 등록",   "강의계획서 업로드 → Step 0"),
    ("05", "데모 — 노트 생성",   "원클릭 녹음 → 실시간 스트리밍"),
    ("06", "다음 단계",           "10주차 베타 테스트 계획"),
]

cols = [0.4, 6.9]
rows = [1.25, 2.9, 4.55]

for i, (num, title, sub) in enumerate(toc):
    cx = cols[i % 2]
    cy = rows[i // 2]
    rect(s, cx, cy, 6.1, 1.45, fill=NAVY_CARD)
    oval(s, cx + 0.18, cy + 0.38, 0.55, 0.55, fill=SKY_BLUE)
    tb(s, cx + 0.18, cy + 0.38, 0.55, 0.55, num, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, cx + 0.85, cy + 0.18, 5.0, 0.52, title, size=18, color=WHITE, bold=True)
    tb(s, cx + 0.85, cy + 0.72, 5.0, 0.45, sub, size=13, color=GRAY)

footer_skyblue(s)
notes(s, """오늘 발표는 크게 세 파트로 구성됩니다.
첫 번째로 현재 개발 진행 상황과 로드맵 업데이트를 말씀드리겠습니다.
두 번째로 이번에 완성된 알파 빌드의 주요 기능을 실제 화면 데모로 보여드리겠습니다.
마지막으로 다음 단계인 10주차 베타 테스트 계획을 말씀드리겠습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 3 — SDLC 진행 현황
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "SDLC 진행 현황", "구현(Implementation) 단계 — 75% 완료")

bars = [
    ("요구분석",  100, GREEN,      "✓ 완료 (3주차)"),
    ("시스템설계", 100, GREEN,     "✓ 완료 (4주차)"),
    ("구   현",    75, ORANGE,     "◀ 현재 — 알파 빌드 완성"),
    ("테스트",      0, LIGHT_GRAY, "예정 (10주차)"),
    ("배   포",     0, LIGHT_GRAY, "예정 (13주차)"),
]
bar_x, bar_w = 2.1, 7.2
by = 1.55
for label, pct, color, note in bars:
    tb(s, 0.3, by, 1.75, 0.38, label, size=14, color=CHARCOAL, bold=(pct > 0))
    rect(s, bar_x, by + 0.06, bar_w, 0.28, fill=LIGHT_GRAY)
    if pct > 0:
        rect(s, bar_x, by + 0.06, bar_w * pct / 100, 0.28, fill=color)
    pct_col = GREEN if color == GREEN else (ORANGE if pct > 0 else GRAY)
    tb(s, bar_x + bar_w + 0.15, by, 0.75, 0.38, f"{pct}%", size=13, color=pct_col, bold=True)
    tb(s, bar_x + bar_w + 0.95, by, 4.0, 0.38, note, size=12, color=GRAY)
    by += 0.58

tb(s, 0.4, 4.65, 12, 0.42, "구현 75% 세부 근거", size=16, color=NAVY, bold=True)

details = [
    ("백엔드 API 서버 (FastAPI, 엔드포인트 9개)", "✅ 완료", True),
    ("데이터베이스 스키마 (SQLite + SQLAlchemy)", "✅ 완료", True),
    ("PDF 파서 / STT / 데이터 병합", "✅ 완료", True),
    ("LLM 연동 + 프롬프트 + SSE 스트리밍", "✅ 완료", True),
    ("프론트엔드 UI (모바일 앱 스타일)", "✅ 완료  ← 신규", True),
    ("서버 배포 / 인프라", "⬜ 미착수", False),
    ("외부 사용자 테스트", "⬜ 미착수", False),
]
dx, dw = 0.4, 6.0
for i, (task, status, done) in enumerate(details):
    col, row = i % 2, i // 2
    x = dx + col * 6.45
    y = 5.1 + row * 0.38
    bg = GREEN_BG if done else LIGHT_GRAY
    sc = GREEN if done else GRAY
    rect(s, x, y, dw, 0.33, fill=bg)
    tb(s, x + 0.12, y + 0.02, 4.4, 0.29, task, size=11, color=CHARCOAL)
    tb(s, x + 4.55, y + 0.02, 1.4, 0.29, status, size=11, color=sc, bold=True)

notes(s, """저희는 폭포수+애자일 하이브리드 방식으로 개발하고 있습니다.
3차 발표에서 구현 단계 35%라고 말씀드렸는데, 이번에 75%까지 완료했습니다.
가장 큰 변화는 3차 발표에서 미착수였던 프론트엔드 UI가 이번에 완성됐다는 점입니다.
SSE 실시간 스트리밍, 인앱 녹음, 노트 편집/삭제까지 모두 동작 확인이 끝났습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 4 — 프로젝트 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "프로젝트 로드맵", "4차 발표 마일스톤 달성")

# 발표 라벨 박스
labels = [
    (0.7,  "1차 발표", SKY_BLUE),
    (3.0,  "2차 발표", SKY_BLUE),
    (5.3,  "3차 발표", SKY_BLUE),
    (7.6,  "4차 발표 ★", ORANGE),
]
for lx, ltxt, lcolor in labels:
    rect(s, lx, 1.55, 1.95, 0.42, fill=lcolor)
    tb(s, lx, 1.55, 1.95, 0.42, ltxt, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 주차 라벨
weeks = ["1-2주차", "3주차", "4주차", "5주차", "8주차", "10주차", "13주차"]
wx = [0.5, 2.3, 4.1, 6.0, 7.8, 10.0, 12.0]
for i, (w, x) in enumerate(zip(weeks, wx)):
    tb(s, x - 0.3, 2.15, 1.8, 0.38, w, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# 타임라인 선
rect(s, 0.8, 2.78, 11.7, 0.05, fill=GRAY)

# 마일스톤 원
milestones = [
    (0.8,  GREEN,    True,   False),  # 1-2주차 완료
    (2.6,  GREEN,    True,   False),  # 3주차
    (4.4,  GREEN,    True,   False),  # 4주차
    (6.3,  GREEN,    True,   False),  # 5주차
    (8.1,  ORANGE,   True,   True),   # 8주차 현재
    (10.3, LIGHT_GRAY, False, False), # 10주차
    (12.3, LIGHT_GRAY, False, False), # 13주차
]
for mx, mc, filled, current in milestones:
    r = 0.32
    if filled:
        oval(s, mx - r/2, 2.62, r, r, fill=mc)
    else:
        shape = s.shapes.add_shape(9, Inches(mx - r/2), Inches(2.62), Inches(r), Inches(r))
        shape.fill.background()
        shape.line.color.rgb = GRAY
        from pptx.util import Pt as Pt2
        shape.line.width = Pt2(1.5)
    if current:
        tb(s, mx - 0.5, 2.2, 1.0, 0.32, "★", size=18, color=ORANGE, align=PP_ALIGN.CENTER)

# 하단 레이블
sublabels = [
    (0.5,  "아이디어\n기획"),
    (2.3,  "요구분석\n완료"),
    (4.1,  "시스템설계\n완료"),
    (6.0,  "백엔드MVP\n완료"),
    (7.8,  "알파빌드\n완료 ★"),
    (10.0, "베타테스트\n피드백"),
    (12.0, "최종발표\n런칭로드맵"),
]
for sx, stxt in sublabels:
    tb(s, sx - 0.35, 3.05, 2.0, 0.72, stxt, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# 달성 항목
tb(s, 0.5, 3.95, 12, 0.42, "이번 마일스톤 달성 항목", size=16, color=SKY_BLUE, bold=True)
achievements = [
    "웹 프론트엔드 개발 완성 — 사용자가 코드 없이 사용 가능한 UI",
    "브라우저 인앱 녹음 기능 통합 — 원클릭 시작/중지",
    "SSE 실시간 스트리밍 — 노트가 타이핑되듯 실시간 표시",
    "노트 조회·편집·삭제·다운로드 전체 기능 완성",
]
for i, ach in enumerate(achievements):
    y = 4.45 + i * 0.48
    oval(s, 0.5, y + 0.06, 0.28, 0.28, fill=GREEN)
    tb(s, 0.9, y, 11.5, 0.42, ach, size=14, color=WHITE)

footer_navy(s, "3차 발표 예고 → 4차 발표 달성 완료")

notes(s, """화면의 로드맵에서 별표 표시가 현재 저희 위치입니다. 8주차, 4차 발표 시점입니다.
3차 발표에서 8주차 목표로 설정한 네 가지 항목을 모두 달성했습니다.
특히 원클릭 녹음과 실시간 스트리밍은 저희 서비스의 핵심 차별화 기능입니다.
다음 목표는 10주차 베타 테스트로, 외부 사용자에게 실제로 써보게 하고 피드백을 받겠습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 5 — 알파 빌드 주요 기능
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "알파 빌드 주요 기능", "사용자가 코드 없이 쓸 수 있는 완성 빌드")

features = [
    (SKY_BLUE, "01", "모바일 앱 스타일 UI",
     "최대 430px 앱 셸\n하단 네비게이션 바\n터치 최적화 버튼 크기\n카드 기반 레이아웃"),
    (GREEN, "02", "원클릭 강의 등록",
     "강의계획서 PDF 업로드\nGemini가 커리큘럼 전체 분석\n과목명·교수명 자동 추출\n주차별 Step 0 로드맵 즉시 확인"),
    (ORANGE, "03", "실시간 스트리밍 노트 생성",
     "음성·PDF·필기 동시 업로드\n인앱 원클릭 녹음 지원\n노트가 타이핑되듯 실시간 표시\n완료 시 브라우저 알림"),
]

for i, (color, num, title, body) in enumerate(features):
    x = 0.5 + i * 4.25
    # 카드 배경
    rect(s, x, 1.55, 4.05, 4.85, fill=NAVY_CARD)
    # 상단 색상 바
    rect(s, x, 1.55, 4.05, 0.12, fill=color)
    # 번호 원
    oval(s, x + 1.65, 1.75, 0.7, 0.7, fill=color)
    tb(s, x + 1.65, 1.75, 0.7, 0.7, num, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 제목
    tb(s, x + 0.2, 2.58, 3.65, 0.58, title, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 구분선
    rect(s, x + 0.3, 3.22, 3.45, 0.04, fill=color)
    # 본문
    tb(s, x + 0.25, 3.35, 3.6, 2.8, body, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 추가 기능 행
tb(s, 0.5, 6.52, 12, 0.32, "추가 기능: 노트 인라인 편집 · 개별 삭제 · HTML 다운로드 · 탭 자동 연결 · 완료 알림", size=13, color=GRAY, align=PP_ALIGN.CENTER)

footer_navy(s, "3가지 핵심 기능 + 9개 API 엔드포인트 완성")

notes(s, """이번 알파 빌드에서 완성된 세 가지 핵심 기능을 소개하겠습니다.
첫 번째는 모바일 앱 스타일 UI입니다. 스마트폰처럼 보이는 최대 430px 앱 셸에 하단 네비게이션 바를 적용해서, PC에서도 앱처럼 사용할 수 있습니다.
두 번째는 원클릭 강의 등록입니다. 강의계획서를 업로드하면 과목명과 교수명을 자동으로 추출하고 전체 커리큘럼을 분석해서 Step 0 로드맵을 만들어줍니다.
세 번째는 실시간 스트리밍입니다. 기존에는 노트 생성에 1~2분을 기다려야 했는데, 이제는 첫 문장이 나오는 순간부터 실시간으로 화면에 표시됩니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 6 — 데모: 강의 등록
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "데모 ① — 강의 등록", "강의계획서 PDF 업로드 → AI 자동 분석")

imgs_row1 = [
    ("01_강의계획서 등록.png",     0.3,  1.5, 3.8),
    ("02_강의계획서 등록중.png",   4.3,  1.5, 3.8),
    ("03_강의계획서 등록완료.png", 8.3, 1.5, 3.8),
]
captions1 = [
    (0.3,  "① 강의계획서 PDF 드래그앤드롭"),
    (4.3,  "② 업로드 중 — 로딩 표시"),
    (8.3,  "③ 등록 완료 — 과목명·교수명 자동 추출"),
]
for fname, x, y, w in imgs_row1:
    pic(s, os.path.join(IMG_DIR, fname), x, y, w=w)
for x, cap in captions1:
    tb(s, x, 5.35, 3.8, 0.45, cap, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Step 0 분석 결과
rect(s, 0.3, 5.9, 12.7, 1.05, fill=GREEN_BG)
tb(s, 0.5, 5.97, 4.5, 0.42, "Step 0 출력 (강의 로드맵)", size=14, color=NAVY, bold=True)
tb(s, 0.5, 6.38, 12.3, 0.42,
   "과목 전체 흐름 · 주차별 핵심 개념 · 시험 출제 가능 포인트 자동 분석 → 이후 매주 노트 생성 시 맥락으로 자동 주입",
   size=12, color=CHARCOAL)

notes(s, """강의 등록은 학기 초 한 번만 하는 과정입니다.
강의계획서 PDF를 드래그앤드롭하면 Gemini가 전체 커리큘럼을 읽고 과목명과 교수명을 자동으로 추출합니다.
이어서 주차별 핵심 개념, 시험 출제 포인트 등 Step 0 로드맵이 생성됩니다.
이 로드맵은 이후 매주 노트를 생성할 때 자동으로 맥락으로 주입됩니다. 사용자가 별도로 프롬프트를 작성할 필요가 없습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 7 — 데모: 강의계획서 분석 결과
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "데모 ② — Step 0 강의 분석 결과", "교수님 커리큘럼 기반 맥락 자동 생성")

pic(s, os.path.join(IMG_DIR, "04_강의계획서 분석.png"), 0.4, 1.4, h=5.3)

# 우측 설명 카드
card_x = 6.0
cards = [
    (SKY_BLUE,  "전체 강의 흐름 파악",    "주차별 주제와 흐름을 한눈에\n시험까지 전체 맥락 이해"),
    (GREEN,     "핵심 개념 사전 식별",    "어떤 개념이 중요한지\n미리 파악 → 노트 중요도 반영"),
    (ORANGE,    "시험 포인트 예측",       "출제 빈도 높은 단원 표시\n셀프 시험 문제 생성 기준"),
]
cy = 1.4
for color, title, body in cards:
    rect(s, card_x, cy, 6.9, 1.6, fill=LIGHT_GRAY)
    rect(s, card_x, cy, 0.12, 1.6, fill=color)
    tb(s, card_x + 0.25, cy + 0.18, 6.4, 0.5, title, size=16, color=NAVY, bold=True)
    tb(s, card_x + 0.25, cy + 0.72, 6.4, 0.72, body, size=13, color=GRAY)
    cy += 1.75

notes(s, """화면은 강의계획서 분석 결과 화면입니다.
Gemini가 강의계획서를 읽고 세 가지 정보를 자동으로 생성합니다.
첫째, 전체 강의 흐름입니다. 주차별 주제와 연결 구조를 파악합니다.
둘째, 핵심 개념 목록입니다. 이 정보는 이후 노트 생성 시 중요도 별점에 반영됩니다.
셋째, 시험 포인트 예측입니다. 이것이 셀프 시험 문제 생성의 기준이 됩니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 8 — 데모: 노트 생성
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "데모 ③ — 노트 생성", "원클릭 녹음 → STT → 실시간 스트리밍 노트")

imgs = [
    ("05_강의 선택.png",    0.3, 1.5, 2.9),
    ("06_강의 자료 입력.png", 3.4, 1.5, 2.9),
    ("07_음성변환.png",     6.5, 1.5, 2.9),
    ("08_노트 생성 완료.png", 9.6, 1.5, 3.4),
]
caps = [
    (0.3,  "① 강의 선택"),
    (3.4,  "② 자료·녹음 입력"),
    (6.5,  "③ 음성 → 텍스트 변환"),
    (9.6,  "④ 실시간 스트리밍 완료"),
]
for fname, x, y, w in imgs:
    pic(s, os.path.join(IMG_DIR, fname), x, y, w=w)
for x, cap in caps:
    tb(s, x, 5.3, 3.0, 0.42, cap, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# SSE 설명 박스
rect(s, 0.3, 5.82, 12.7, 1.12, fill=ORANGE_BG)
rect(s, 0.3, 5.82, 0.1, 1.12, fill=ORANGE)
tb(s, 0.55, 5.9, 5.5, 0.45, "핵심 UX — SSE 실시간 스트리밍", size=14, color=CHARCOAL, bold=True)
tb(s, 0.55, 6.36, 12.2, 0.45,
   "기존: 1~2분 기다린 뒤 노트 일괄 표시  →  이제: 첫 문장부터 타이핑되듯 실시간 표시 (완료 시 브라우저 알림)",
   size=12, color=CHARCOAL)

notes(s, """노트 생성 과정을 단계별로 보여드리겠습니다.
첫 번째로 강의를 선택합니다. 등록된 강의 카드 목록에서 해당 강의를 탭합니다.
두 번째로 자료를 입력합니다. 수업 PDF를 올리고, 인앱 녹음 버튼을 누르거나 파일을 업로드합니다. 주차 번호도 입력합니다.
세 번째로 음성이 변환됩니다. Gemini 멀티모달이 음성 파일을 한국어 텍스트로 받아씁니다.
네 번째가 핵심입니다. 노트가 타이핑되듯 실시간으로 화면에 표시되고, 완료되면 브라우저 알림이 옵니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 9 — 데모: 노트 조회/편집
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "데모 ④ — 노트 조회 · 편집", "생성된 노트 조회 · 인라인 편집 · 다운로드")

imgs2 = [
    ("09_노트 조회.png",  0.3, 1.5, 3.1),
    ("10_노트 편집.png",  3.7, 1.5, 3.1),
    ("11_복습문제 예시.png", 7.1, 1.5, 6.0),
]
caps2 = [
    (0.3, "⑤ 노트 조회 — 주차별 목록"),
    (3.7, "⑥ 인라인 편집 · 저장"),
]
for fname, x, y, w in imgs2:
    pic(s, os.path.join(IMG_DIR, fname), x, y, w=w)
for x, cap in caps2:
    tb(s, x, 5.35, 3.2, 0.42, cap, size=12, color=GRAY, align=PP_ALIGN.CENTER)

tb(s, 7.1, 5.35, 6.0, 0.42, "⑦ 5섹션 노트 + 셀프 시험 문제", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# 기능 요약 바
rect(s, 0.3, 5.9, 12.7, 1.12, fill=GREEN_BG)
rect(s, 0.3, 5.9, 0.1, 1.12, fill=GREEN)
tb(s, 0.55, 5.98, 5.5, 0.42, "노트 관리 기능 요약", size=14, color=CHARCOAL, bold=True)
tb(s, 0.55, 6.42, 12.2, 0.45,
   "조회 · 인라인 편집/저장 · 개별 삭제 · HTML 다운로드 (브라우저 Ctrl+P → PDF) · 완료 후 탭 자동 이동",
   size=12, color=CHARCOAL)

notes(s, """생성된 노트는 세 가지 방법으로 관리할 수 있습니다.
첫 번째로 노트 조회입니다. 강의별로 주차 노트 목록을 확인하고, 클릭하면 5섹션으로 구성된 노트 내용이 펼쳐집니다.
두 번째로 인라인 편집입니다. 편집 버튼을 누르면 textarea로 전환되어 앱 안에서 바로 수정하고 저장할 수 있습니다.
세 번째로 노트 맨 아래 셀프 시험 섹션입니다. 서술형 2문제와 객관식 3문제가 강의계획서 기반으로 자동 생성됩니다. 시험 직전에 복습 도구로 쓸 수 있습니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 10 — 기술 스택 업데이트
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "기술 스택 업데이트", "3차 발표 대비 프론트엔드 계층 추가")

stacks = [
    ("AI",       SKY_BLUE, "Gemini 2.5-flash",
     "LLM 노트 생성 · STT 음성 변환 · 강의계획서 분석"),
    ("백엔드",   GREEN,    "FastAPI + Python",
     "REST API · SSE 스트리밍 · 9개 엔드포인트"),
    ("데이터",   ORANGE,   "SQLite + SQLAlchemy",
     "강의 · 노트 · 소스 3테이블 구조"),
    ("PDF 처리", SKY_BLUE, "PyMuPDF (fitz)",
     "PDF → 텍스트 추출"),
    ("프론트엔드 ★ 신규", GREEN, "Vanilla HTML/CSS/JS",
     "단일 파일 SPA · 빌드 도구 없음 · 모바일 430px 셸"),
    ("마크다운 ★ 신규", ORANGE, "marked.js (CDN)",
     "생성 노트 실시간 렌더링"),
    ("녹음 ★ 신규", SKY_BLUE, "MediaRecorder API",
     "브라우저 내 원클릭 음성 녹음"),
]

col1_x, col2_x = 0.4, 6.9
for i, (layer, color, tech, desc) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = col1_x if col == 0 else col2_x
    y = 1.5 + row * 1.18
    rect(s, x, y, 6.1, 1.05, fill=NAVY_CARD)
    rect(s, x, y, 0.1, 1.05, fill=color)
    tb(s, x + 0.25, y + 0.08, 2.5, 0.38, layer, size=12, color=color, bold=True)
    tb(s, x + 2.75, y + 0.08, 3.2, 0.38, tech, size=13, color=WHITE, bold=True)
    tb(s, x + 0.25, y + 0.58, 5.7, 0.38, desc, size=12, color=GRAY)

footer_navy(s, "프론트엔드 추가로 풀스택 구성 완료")

notes(s, """3차 발표에서는 백엔드와 AI 계층만 있었다면, 이번에 프론트엔드 계층이 추가됐습니다.
특히 빌드 도구 없이 단일 HTML 파일로 구성된 SPA는 서버에 별도 정적 파일 서빙이 필요 없어서, 파일을 브라우저에서 열기만 하면 바로 사용 가능합니다.
MediaRecorder API를 활용한 원클릭 녹음은 별도 앱 설치 없이 브라우저에서 바로 작동합니다.""")

# ════════════════════════════════════════════════════════════
# SLIDE 11 — 다음 단계
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "다음 단계", "베타 테스트 → 최종 발표")

next_steps = [
    (SKY_BLUE, "10주차", "베타 테스트",
     ["외부 사용자 5~10명 모집",
      "실제 수업에 적용 · 피드백 수집",
      "사용성 개선 및 오류 수정",
      "Whisper 오픈소스 전환 검토"]),
    (GREEN, "12주차", "고도화",
     ["사진(칠판 필기) 업로드 — Gemini OCR",
      "노트 키워드 검색 기능",
      "PWA 패키징 — 실제 앱 아이콘",
      "비용 구조 재검증 (역마진 해소)"]),
    (ORANGE, "13주차", "최종 발표",
     ["서비스 런칭 로드맵 발표",
      "누적 테스트 사용자 데이터",
      "수익 구조 최종 검증",
      "시장 진출 전략"]),
]

for i, (color, week, title, items) in enumerate(next_steps):
    x = 0.5 + i * 4.25
    rect(s, x, 1.5, 4.05, 5.1, fill=NAVY_CARD)
    rect(s, x, 1.5, 4.05, 0.12, fill=color)
    oval(s, x + 1.65, 1.7, 0.7, 0.7, fill=color)
    tb(s, x + 1.65, 1.7, 0.7, 0.7, week, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + 0.2, 2.52, 3.65, 0.52, title, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    rect(s, x + 0.3, 3.1, 3.45, 0.04, fill=color)
    for j, item in enumerate(items):
        tb(s, x + 0.3, 3.22 + j * 0.52, 3.65, 0.48, f"• {item}", size=12, color=LIGHT_GRAY)

footer_skyblue(s, "이상으로 렉처노트 4차 발표를 마치겠습니다.  |  경영전략 9팀")

notes(s, """이번 알파 빌드 완성으로 저희는 '코드로 증명'하는 단계를 마쳤습니다.
다음 단계인 10주차에는 외부 사용자 베타 테스트를 진행합니다. 실제 학생들이 수업에 써보고 피드백을 받겠습니다.
12주차에는 칠판 사진 OCR, PWA 패키징 등 고도화 기능을 추가합니다.
13주차 최종 발표에서는 서비스 런칭 로드맵과 수익 구조 최종안을 발표하겠습니다.
이상으로 렉처노트 4차 발표를 마치겠습니다. 감사합니다.""")

# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"\n[OK] 저장 완료: {OUTPUT_PATH}")
print(f"   슬라이드 수: {len(prs.slides)}장")
