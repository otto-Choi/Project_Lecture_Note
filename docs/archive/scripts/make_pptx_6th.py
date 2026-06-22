"""
렉처노트 6차(최종) 발표 PPTX 생성 스크립트
design.md 디자인 시스템 적용
대본 없이 슬라이드만으로 발표 가능하도록 정보를 충분히 담음
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_DIR = r"g:\내 드라이브\02_Study\01_Univ\2026학년도 1학기\06_경영전략\조별\lecture_note_project"
IMG_DIR  = os.path.join(BASE_DIR, "docs", "presentations", "5차", "screenshots")  # 5차 스크린샷 재활용
OUTPUT_PATH = os.path.join(
    BASE_DIR, "docs", "presentations", "6차", "렉처노트_6차발표.pptx"
)

os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)

# ─── 컬러 팔레트 ─────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x2D, 0x5E)
NAVY_CARD  = RGBColor(0x22, 0x38, 0x72)
SKY_BLUE   = RGBColor(0x4B, 0x82, 0xC5)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
GREEN_BG   = RGBColor(0xE8, 0xF8, 0xEE)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_BG  = RGBColor(0xFE, 0xF3, 0xC7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF7)
CHARCOAL   = RGBColor(0x1C, 0x1C, 0x1E)
GRAY       = RGBColor(0x6E, 0x6E, 0x73)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


# ─── 헬퍼 함수 ───────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, line_pt=0.75):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def oval(slide, x, y, w, h, fill=None):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def tb(slide, x, y, w, h, text, size=16, color=WHITE, bold=False, italic=False,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.color.rgb  = color
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.name       = "맑은 고딕"
    return box

def tb2(slide, x, y, w, h, lines, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, line_spacing=1.15):
    """여러 줄 텍스트박스 (줄바꿈 포함)"""
    from pptx.util import Pt as Pt2
    from pptx.oxml.ns import qn
    from lxml import etree
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size      = Pt2(size)
        run.font.color.rgb = color
        run.font.bold      = bold
        run.font.name      = "맑은 고딕"
    return box

def pic(slide, img_path, x, y, w=None, h=None):
    if not os.path.exists(img_path):
        print(f"  [경고] 이미지 없음: {img_path}")
        return None
    if w and h:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
    elif w:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(img_path, Inches(x), Inches(y), height=Inches(h))

def navy_bg(slide):
    rect(slide, 0, 0, 13.333, 7.5, fill=NAVY)

def footer_sky(slide, text="경영전략 팀 프로젝트  |  9팀 6차 발표  |  2026.06"):
    rect(slide, 0, 6.92, 13.333, 0.58, fill=SKY_BLUE)
    tb(slide, 0, 6.93, 13.333, 0.56, text, size=13, color=WHITE, align=PP_ALIGN.CENTER)

def footer_navy_bar(slide, text):
    rect(slide, 0, 6.92, 13.333, 0.58, fill=NAVY_CARD)
    tb(slide, 0.3, 6.93, 12.7, 0.56, text, size=13, color=WHITE, align=PP_ALIGN.CENTER)

def slide_title_white(slide, title, subtitle=None):
    tb(slide, 0.5, 0.2, 12.3, 0.72, title, size=30, color=NAVY, bold=True)
    if subtitle:
        rect(slide, 0.5, 0.95, 12.3, 0.03, fill=SKY_BLUE)
        tb(slide, 0.5, 1.0, 12.3, 0.42, subtitle, size=16, color=GRAY)

def slide_title_navy(slide, title, subtitle=None):
    tb(slide, 0.5, 0.2, 12.3, 0.72, title, size=30, color=WHITE, bold=True)
    if subtitle:
        rect(slide, 0.5, 0.95, 8.0, 0.03, fill=SKY_BLUE)
        tb(slide, 0.5, 1.0, 12.3, 0.42, subtitle, size=16, color=SKY_BLUE)

def badge(slide, x, y, text, color):
    rect(slide, x, y, 1.1, 0.3, fill=color)
    tb(slide, x, y, 1.1, 0.3, text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)

rect(s, 0.85, 1.3, 0.07, 4.5, fill=SKY_BLUE)
tb(s, 1.05, 1.15, 10, 0.45, "경영전략 9팀  |  6차 발표 (최종)", size=16, color=GRAY)
tb(s, 1.05, 1.65, 11, 1.3, "LectureNote", size=78, color=WHITE, bold=True)
tb(s, 1.05, 2.85, 10, 0.65, "렉쳐노트", size=38, color=SKY_BLUE, bold=True)
rect(s, 1.05, 3.60, 5.5, 0.06, fill=SKY_BLUE)
tb(s, 1.05, 3.74, 11.5, 0.55,
   "앱 출시 완성 — Android 패키징 · Google Play 등록 · 베타 확장",
   size=20, color=WHITE, bold=True)
tb(s, 1.05, 4.36, 11.5, 0.45,
   "5차 발표에서 예고한 앱 패키징 완료 | 13주차 최종 발표",
   size=15, color=GRAY, italic=True)

# 핵심 수치 3개
nums = [
    ("v0.4.0", "현재 버전"),
    ("19개", "API 엔드포인트"),
    ("Play Store", "등록 검토 중"),
]
for i, (n, l) in enumerate(nums):
    bx = 1.05 + i * 3.5
    rect(s, bx, 5.05, 3.1, 0.88, fill=NAVY_CARD)
    tb(s, bx, 5.1, 3.1, 0.48, n, size=22, color=SKY_BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, bx, 5.6, 3.1, 0.28, l, size=12, color=GRAY, align=PP_ALIGN.CENTER)

tb(s, 1.05, 6.05, 11, 0.42,
   "이병헌 (팀장) · 김도환 · 최철원  |  2026년 6월",
   size=15, color=GRAY)
footer_sky(s)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "발표 구성")

toc = [
    ("01", "SDLC 진행 현황",        "구현 100% · 테스트·배포 단계 완료"),
    ("02", "프로젝트 로드맵",        "1~6차 전체 마일스톤 달성 현황"),
    ("03", "6차 신규 완료 항목",     "패키징 · Play Store · 베타 확장 · 피드백 반영"),
    ("04", "Android 앱 패키징",     "Capacitor 6 래핑 구조 및 설치 방법"),
    ("05", "베타 테스트 & 피드백",  "확장된 인터뷰 결과 및 반영 현황"),
    ("06", "최종 기술 스택",         "전체 아키텍처 및 배포 구성"),
    ("07", "런칭 로드맵 & 수익 구조", "Phase 1~3 · 구독 모델 최종안"),
]

cols = [0.4, 6.9]
rows = [1.3, 2.72, 4.14, 5.56]

for i, (num, title, sub) in enumerate(toc):
    cx = cols[i % 2]
    cy = rows[i // 2]
    rect(s, cx, cy, 6.1, 1.22, fill=NAVY_CARD)
    oval(s, cx + 0.18, cy + 0.32, 0.52, 0.52, fill=SKY_BLUE)
    tb(s, cx + 0.18, cy + 0.32, 0.52, 0.52, num, size=12, color=WHITE,
       bold=True, align=PP_ALIGN.CENTER)
    tb(s, cx + 0.85, cy + 0.14, 5.0, 0.46, title, size=17, color=WHITE, bold=True)
    tb(s, cx + 0.85, cy + 0.62, 5.0, 0.42, sub, size=12, color=GRAY)

footer_sky(s)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — SDLC 진행 현황
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "SDLC 진행 현황",
                  "5차 대비: 테스트·배포 단계 완료 — 앱 출시 가능 상태")

bars = [
    ("요구분석",   100, GREEN,      "✓ 완료 (3주차)"),
    ("시스템설계", 100, GREEN,      "✓ 완료 (4주차)"),
    ("구   현",   100, GREEN,      "✓ 완료 (10주차) — API 19개 · 화면 10개+"),
    ("테스트",      95, SKY_BLUE,  "◀ 95% — 베타 테스트 확장 완료"),
    ("배   포",     90, ORANGE,    "◀ 90% — Render LIVE · Play Store 검토 중"),
]
bar_x, bar_w = 2.1, 7.4
by = 1.55
for label, pct, color, note in bars:
    tb(s, 0.3, by, 1.75, 0.38, label, size=14, color=CHARCOAL, bold=(pct == 100))
    rect(s, bar_x, by + 0.06, bar_w, 0.28, fill=LIGHT_GRAY)
    if pct > 0:
        rect(s, bar_x, by + 0.06, bar_w * pct / 100, 0.28, fill=color)
    pct_col = GREEN if color == GREEN else (SKY_BLUE if color == SKY_BLUE else ORANGE)
    tb(s, bar_x + bar_w + 0.15, by, 0.75, 0.38, f"{pct}%", size=13, color=pct_col, bold=True)
    tb(s, bar_x + bar_w + 0.95, by, 4.2, 0.38, note, size=12, color=GRAY)
    by += 0.58

# 6차 신규 완료 박스
rect(s, 0.4, 4.6, 12.5, 2.2, fill=GREEN_BG)
rect(s, 0.4, 4.6, 0.1, 2.2, fill=GREEN)
tb(s, 0.65, 4.68, 8, 0.4, "6차 발표 신규 완료 항목", size=15, color=NAVY, bold=True)

items6 = [
    ("Android APK 패키징",         "Capacitor 6 래핑 · 아이콘/스플래시/상태바 완성"),
    ("Google Play Store 등록",     "appId: com.kbslab.lecturenote · 검토 예정"),
    ("베타 테스트 확장",            "인터뷰 확장 · 실제 수업 적용 데이터 수집"),
    ("피드백 반영 완료",            "다크모드 개선 · 키워드 검색 · Gen-modal 미리보기 · SVG 안정화"),
]
for i, (title, desc) in enumerate(items6):
    col = i % 2
    row = i // 2
    x = 0.65 + col * 6.2
    y = 5.12 + row * 0.55
    tb(s, x, y, 1.5, 0.38, f"✓ {title}", size=12, color=GREEN, bold=True)
    tb(s, x + 1.55, y, 4.5, 0.38, desc, size=12, color=CHARCOAL)

footer_navy_bar(s, "5차: 테스트 80% · 배포 80%  →  6차: 테스트 95% · 배포 90% (Play Store 검토)")


# ════════════════════════════════════════════════════════════
# SLIDE 4 — 프로젝트 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "프로젝트 로드맵", "6차 발표 — 최종 마일스톤 달성")

labels = [
    (0.5,  "1차 발표", SKY_BLUE),
    (2.5,  "2차 발표", SKY_BLUE),
    (4.5,  "3차 발표", SKY_BLUE),
    (6.5,  "4차 발표", SKY_BLUE),
    (8.5,  "5차 발표", SKY_BLUE),
    (10.5, "6차 발표 ★", ORANGE),
]
for lx, ltxt, lcolor in labels:
    rect(s, lx, 1.5, 2.05, 0.38, fill=lcolor)
    tb(s, lx, 1.5, 2.05, 0.38, ltxt, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

weeks  = ["1-2주", "3-4주", "5-6주", "8주", "10주", "13주"]
wx     = [0.7, 2.7, 4.7, 6.7, 8.7, 10.7]
for w, x in zip(weeks, wx):
    tb(s, x - 0.2, 2.05, 1.8, 0.32, w, size=11, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.85, 2.65, 11.6, 0.05, fill=GRAY)

milestones = [
    (0.85,  GREEN,      "아이디어\n기획"),
    (2.85,  GREEN,      "설계\n완료"),
    (4.85,  GREEN,      "백엔드\nMVP"),
    (6.85,  GREEN,      "알파\n빌드"),
    (8.85,  GREEN,      "베타\n빌드"),
    (10.85, ORANGE,     "앱 출시\n완성 ★"),
]
for mx, mc, lbl in milestones:
    r = 0.32
    oval(s, mx - r/2, 2.49, r, r, fill=mc)
    tb(s, mx - 0.55, 2.92, 1.6, 0.65, lbl, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# 달성 항목
tb(s, 0.5, 3.75, 12.3, 0.4, "6차 마일스톤 달성 항목", size=15, color=SKY_BLUE, bold=True)
ach6 = [
    "Android APK 패키징 완료 — Capacitor 6 · 아이콘/스플래시/상태바 완성",
    "Google Play Store 등록 (appId: com.kbslab.lecturenote) — 검토 예정",
    "베타 테스트 확장 — 실제 수업 적용 · 인터뷰 확장",
    "피드백 반영 완료 — 다크모드 대비 보정 · 노트 키워드 검색 · Gen-modal 실시간 미리보기 · SVG 아이콘 안정화",
]
for i, ach in enumerate(ach6):
    oval(s, 0.5, 4.22 + i * 0.5, 0.28, 0.28, fill=GREEN)
    tb(s, 0.88, 4.18 + i * 0.5, 11.9, 0.44, ach, size=13, color=WHITE)

footer_navy_bar(s, "1차 아이디어 기획 → 6차 앱 출시 완성  |  13주 완주")


# ════════════════════════════════════════════════════════════
# SLIDE 5 — 6차 신규 완료 항목 (4카드)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "6차 신규 완료 항목", "5차 발표 예고 → 6차 완료 확인")

cards5 = [
    (SKY_BLUE,  "01", "Android APK 패키징",
     ["Capacitor 6 WebView 래핑",
      "appId: com.kbslab.lecturenote",
      "아이콘(1024×1024) · 스플래시(2732×2732)",
      "Deep Navy 상태바 · 안전 영역 처리",
      "Render 백엔드 API 연결",
      "APK 사이드로드 즉시 설치 가능"],
     "NEW"),
    (GREEN,     "02", "Google Play Store 등록",
     ["Release AAB (bundleRelease) 생성 완료",
      "개발자 계정 등록 및 앱 정보 입력",
      "스크린샷 · 설명 · 카테고리 설정",
      "콘텐츠 등급 설문 완료",
      "검토 요청 완료 — 승인 대기 중",
      "승인 후 링크 공유 예정"],
     "LIVE"),
    (ORANGE,    "03", "베타 테스트 확장",
     ["5차: 외부 5명 정성 인터뷰",
      "6차: 실제 수업 적용 확장",
      "APK 배포로 설치 접근성 개선",
      "강의 등록 1회 → 매주 노트 생성 패턴 검증",
      "피드백 수집 채널 정비",
      "누적 생성 노트 수 측정"],
     "NEW"),
    (CYAN,      "04", "피드백 반영 완료",
     ["✓ 다크모드 대비 — CSS 변수 전체 재정의",
      "✓ 노트 내 키워드 검색 — 하이라이트·↑↓ 이동",
      "✓ Gen-modal 실시간 미리보기 — 생성 중 텍스트 스트리밍",
      "✓ SVG 아이콘 안정화 — fetch 제거 → 인라인 삽입",
      "✓ 자동 스크롤 — 생성 중 하단 고정·수동 해제",
      "⟳ 복수 파일 업로드 — 백엔드 지원·UI 개선 중"],
     "완료"),
]

CW = 6.1
CH = 5.55
for i, (color, num, title, items, badge_txt) in enumerate(cards5):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.45
    cy = 1.5 + row * 2.92

    rect(s, cx, cy, CW, CH / 2 - 0.02, fill=NAVY_CARD)
    rect(s, cx, cy, CW, 0.1, fill=color)

    # 배지
    badge(s, cx + CW - 1.2, cy + 0.14, badge_txt, color)

    # 번호 원
    oval(s, cx + 0.14, cy + 0.2, 0.46, 0.46, fill=color)
    tb(s, cx + 0.14, cy + 0.2, 0.46, 0.46, num, size=12, color=WHITE,
       bold=True, align=PP_ALIGN.CENTER)

    # 제목
    tb(s, cx + 0.72, cy + 0.18, CW - 1.5, 0.46, title, size=16, color=WHITE, bold=True)

    # 항목 목록
    for j, item in enumerate(items):
        tb(s, cx + 0.18, cy + 0.72 + j * 0.31, CW - 0.3, 0.3, item,
           size=11, color=LIGHT_GRAY)

footer_sky(s)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — Android 앱 패키징 상세
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "Android 앱 패키징 상세",
                  "Capacitor 6 — 기존 웹앱을 Native Android APK로 변환")

# 좌측: 아키텍처 다이어그램
rect(s, 0.4, 1.5, 5.7, 5.28, fill=LIGHT_GRAY)
tb(s, 0.4, 1.5, 5.7, 0.4, "  패키징 아키텍처", size=13, color=CHARCOAL, bold=True)

layers = [
    (GREEN,    "Android APK  (com.kbslab.lecturenote)",      "사용자 기기에 설치되는 네이티브 앱"),
    (SKY_BLUE, "Capacitor 6  WebView 래퍼",                  "Android WebView 위에 웹앱 로드"),
    (ORANGE,   "www/  (HTML·CSS·JS)",                        "기존 Vanilla SPA를 그대로 탑재"),
    (CYAN,     "fetch()  →  Render 백엔드",                  "APK 안에는 프론트만 · 백엔드는 클라우드"),
]
for i, (c, t, d) in enumerate(layers):
    by = 2.0 + i * 1.18
    rect(s, 0.55, by, 5.35, 1.05, fill=c)
    tb(s, 0.7, by + 0.08, 5.1, 0.42, t, size=13, color=WHITE, bold=True)
    tb(s, 0.7, by + 0.54, 5.1, 0.38, d, size=11, color=WHITE)

# 우측: 상세 정보 테이블
right_x = 6.5
tb(s, right_x, 1.5, 6.5, 0.4, "패키징 상세 정보", size=15, color=NAVY, bold=True)

rows_info = [
    ("앱 ID",         "com.kbslab.lecturenote"),
    ("버전",           "0.4.0"),
    ("Capacitor",     "v6.1.2"),
    ("빌드 방식",      "Gradle assembleRelease → AAB"),
    ("Play Store",    "검토 예정 (내부 테스트 트랙)"),
    ("사이드로드",     "APK 파일 직접 설치 가능"),
    ("최소 SDK",       "Android 7.0+ (API 24)"),
    ("권한",           "INTERNET · RECORD_AUDIO · READ_MEDIA"),
    ("상태바",         "Deep Navy (#1A2D5E) 통일"),
    ("아이콘",         "1024×1024 · Adaptive Icon 완성"),
    ("스플래시",       "2732×2732 · Deep Navy 배경"),
    ("API URL",       "앱 내 ⚙ 설정으로 변경 가능"),
]

rect(s, right_x, 1.98, 6.6, 0.3, fill=NAVY)
tb(s, right_x + 0.1, 1.98, 2.0, 0.28, "항목", size=11, color=WHITE, bold=True)
tb(s, right_x + 2.2, 1.98, 4.3, 0.28, "내용", size=11, color=WHITE, bold=True)

for i, (k, v) in enumerate(rows_info):
    by = 2.3 + i * 0.35
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, right_x, by, 6.6, 0.33, fill=bg)
    tc = NAVY if i % 2 == 0 else CHARCOAL
    tb(s, right_x + 0.1, by + 0.02, 2.0, 0.29, k, size=11, color=NAVY, bold=True)
    tb(s, right_x + 2.2, by + 0.02, 4.3, 0.29, v, size=11, color=CHARCOAL)

footer_navy_bar(s, "APK → 사이드로드 즉시 설치 가능  |  Play Store → 승인 후 링크 배포 예정")


# ════════════════════════════════════════════════════════════
# SLIDE 7 — 베타 테스트 & 피드백 반영
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "베타 테스트 확장 & 피드백 반영",
                  "5차(5명 정성 인터뷰) → 6차(APK 배포 기반 확장)")

# 상단 수치 카드 3개
num_cards = [
    (SKY_BLUE, "5명 → 확장",  "인터뷰 규모"),
    (GREEN,    "4건 완료",    "피드백 반영"),
    (ORANGE,   "1건 진행 중", "복수 파일 업로드"),
]
for i, (c, n, l) in enumerate(num_cards):
    bx = 0.4 + i * 4.3
    rect(s, bx, 1.48, 4.0, 0.78, fill=c)
    tb(s, bx, 1.52, 4.0, 0.38, n, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, bx, 1.9, 4.0, 0.28, l, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# 피드백 반영 테이블
rect(s, 0.4, 2.42, 12.5, 0.34, fill=NAVY)
for x_, w_, lbl_ in [(0.5, 1.4, "구분"), (1.95, 3.5, "피드백 내용"),
                     (5.5, 2.5, "5차 상태"), (8.05, 2.0, "6차 상태"), (10.1, 2.7, "처리 방법")]:
    tb(s, x_, 2.44, w_, 0.3, lbl_, size=11, color=WHITE, bold=True)

feedback_rows = [
    ("UX 강점",   "하단 탭 구조 직관적 · 모바일 앱 형태",
     "강점 유지", "강점 유지", "APK에서도 동일 UX 유지"),
    ("기능 강점", "Self-Exam 예상 문제 자동 생성 가장 유용",
     "강점 유지", "강점 유지", "서술형 2 + 객관식 3 포맷 고정"),
    ("UX 명확화", "강의 등록 vs 노트 생성 구분 불명확",
     "진행 중",   "✓ 반영",  "Step 0 배지 + '학기 초 1회' 라벨"),
    ("개선 요청", "다크모드 일부 색 대비 약함",
     "진행 중",   "✓ 반영",  "CSS 변수 전체 재정의 · 배경 오버라이드"),
    ("개선 요청", "수업 자료·녹음 복수 파일 업로드 제한",
     "진행 중",   "⟳ 진행",  "백엔드 list 지원 완료 · UI 개선 중"),
]

for i, (cat, fb, s5, s6, how) in enumerate(feedback_rows):
    by = 2.78 + i * 0.58
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(s, 0.4, by, 12.5, 0.55, fill=bg)

    # 카테고리 배지
    cat_c = GREEN if "강점" in cat else (SKY_BLUE if "명확화" in cat else ORANGE)
    rect(s, 0.45, by + 0.1, 1.35, 0.3, fill=cat_c)
    tb(s, 0.45, by + 0.1, 1.35, 0.3, cat, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, 1.95, by + 0.08, 3.5, 0.38, fb, size=11, color=CHARCOAL)

    s5c = GREEN if "유지" in s5 else (ORANGE if "중" in s5 else CHARCOAL)
    s6c = GREEN if "✓" in s6 else (ORANGE if "⟳" in s6 else CHARCOAL)
    tb(s, 5.5, by + 0.08, 2.5, 0.38, s5, size=11, color=s5c, bold=("유지" in s5))
    tb(s, 8.05, by + 0.08, 2.0, 0.38, s6, size=11, color=s6c, bold=True)
    tb(s, 10.1, by + 0.08, 2.7, 0.38, how, size=10, color=CHARCOAL)

footer_navy_bar(s, "5차 진행 중 2건 → 6차 1건 완료 · 1건 계속 진행")


# ════════════════════════════════════════════════════════════
# SLIDE 8 — 데모: 앱 설치 & 사용 흐름
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "데모 — 앱 설치 및 전체 사용 흐름",
                  "APK 설치 → 회원가입 → 강의 등록 → 노트 생성 → 조회")

# 5단계 흐름 다이어그램 (텍스트 기반)
steps = [
    (SKY_BLUE, "STEP 1",  "APK 설치",
     "APK 파일 전송\n→ 설치 허용\n→ 아이콘 생성"),
    (GREEN,    "STEP 2",  "회원가입 / 로그인",
     "온보딩 화면\n→ 2단계 회원가입\n→ 세션 72h 유지"),
    (ORANGE,   "STEP 3",  "강의 등록\n(학기 초 1회)",
     "강의계획서 PDF\n→ Gemini 분석\n→ Step 0 저장"),
    (CYAN,     "STEP 4",  "노트 생성\n(매주 반복)",
     "녹음 + PDF 입력\n→ 5단계 파이프라인\n→ SSE 스트리밍"),
    (GREEN,    "STEP 5",  "조회 / 편집\n/ 다운로드",
     "주차별 목록\n→ 인라인 편집\n→ HTML · PDF"),
]

total_w = 12.5
step_w  = total_w / len(steps)
sy = 1.5

for i, (c, num, title, body) in enumerate(steps):
    x = 0.4 + i * step_w

    # 연결 화살표
    if i < len(steps) - 1:
        ax = x + step_w - 0.18
        tb(s, ax, sy + 1.5, 0.35, 0.55, "→", size=22, color=GRAY, align=PP_ALIGN.CENTER)

    rect(s, x, sy, step_w - 0.25, 4.8, fill=NAVY_CARD)
    rect(s, x, sy, step_w - 0.25, 0.1, fill=c)

    oval(s, x + (step_w - 0.6) / 2, sy + 0.18, 0.6, 0.6, fill=c)
    tb(s, x + (step_w - 0.6) / 2, sy + 0.18, 0.6, 0.6,
       str(i + 1), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, x + 0.1, sy + 0.9, step_w - 0.45, 0.55, num, size=10, color=c, bold=True)
    tb(s, x + 0.1, sy + 1.4, step_w - 0.45, 0.75, title, size=14, color=WHITE, bold=True)
    rect(s, x + 0.1, sy + 2.2, step_w - 0.5, 0.04, fill=c)
    tb(s, x + 0.1, sy + 2.3, step_w - 0.45, 2.0, body, size=12, color=LIGHT_GRAY)

# 이미지 힌트 박스
rect(s, 0.4, 6.4, 12.5, 0.45, fill=ORANGE_BG)
tb(s, 0.55, 6.42, 12.2, 0.38,
   "▶  실제 화면 데모: 앱 아이콘 → 첫 화면 → 강의 등록 → 노트 생성 → 노트 조회  (라이브 또는 스크린샷)",
   size=12, color=CHARCOAL)

footer_navy_bar(s, "LIVE DEMO  |  lecture-note-2cb6.onrender.com  또는  APK 직접 설치")


# ════════════════════════════════════════════════════════════
# SLIDE 9 — 5차 스크린샷 4장 (데모 보조)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_title_white(s, "앱 화면 스크린샷",
                  "Render 배포 URL 기반 실제 동작 — 모바일 웹 / APK 동일 UI")

imgs_info = [
    ("1_첫화면.png",         0.3,  1.5, 2.95, "① 온보딩 / 로그인"),
    ("5_메인화면.png",       3.55, 1.5, 2.95, "② 강의 등록 메인"),
    ("8_노트생성화면.png",   6.8,  1.5, 2.95, "③ 노트 생성 입력"),
    ("10_노트조회.png",      10.05,1.5, 3.0,  "④ 노트 조회 결과"),
]
for fname, ix, iy, iw, cap in imgs_info:
    fpath = os.path.join(IMG_DIR, fname)
    if os.path.exists(fpath):
        pic(s, fpath, ix, iy, w=iw)
    else:
        rect(s, ix, iy, iw, 4.6, fill=LIGHT_GRAY, line=GRAY)
        tb(s, ix, iy + 2.0, iw, 0.5, f"[{fname}]", size=11, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, ix, 6.15, iw + 0.05, 0.38, cap, size=12, color=GRAY, align=PP_ALIGN.CENTER)

footer_navy_bar(s, "웹 브라우저 · APK · 모바일 3가지 방식 모두 동일 UI 제공")


# ════════════════════════════════════════════════════════════
# SLIDE 10 — 최종 기술 스택
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "최종 기술 스택",
                 "6차 완성 — 인증·배포·패키징 계층까지 풀스택 구성")

# 상단 수치
nums_tech = [
    ("19개", "API 엔드포인트"),
    ("5개",  "DB 테이블"),
    ("10+개","화면 (Screen)"),
    ("3개",  "배포 채널"),
]
for i, (n, l) in enumerate(nums_tech):
    bx = 0.4 + i * 3.2
    rect(s, bx, 1.45, 3.0, 0.7, fill=SKY_BLUE)
    tb(s, bx, 1.48, 3.0, 0.38, n, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, bx, 1.86, 3.0, 0.26, l, size=11, color=WHITE, align=PP_ALIGN.CENTER)

stacks = [
    (SKY_BLUE, "AI",         "Gemini 2.5-flash",
     "LLM 노트 생성 · STT 멀티모달 · 강의계획서 분석 (thinking_budget=0 최적화)"),
    (GREEN,    "백엔드",      "FastAPI + Python",
     "REST API · SSE 스트리밍 · 19개 엔드포인트 · 인라인 스키마 마이그레이션"),
    (ORANGE,   "데이터",      "SQLite + SQLAlchemy",
     "User · Session · Lecture · Source · Output — 5테이블 · cascade 삭제"),
    (CYAN,     "인증",        "SHA256 + salt · 세션 토큰",
     "72h 세션 · 사용자별 데이터 격리 · 자동 로그인 · 비밀번호 변경 시 전체 무효화"),
    (SKY_BLUE, "PDF 처리",   "PyMuPDF (fitz)",
     "PDF → 페이지별 텍스트 추출 · 강의계획서 / 수업자료 공용"),
    (GREEN,    "프론트엔드",  "Vanilla HTML/CSS/JS",
     "단일 파일 SPA · 모바일 430px · SSE 스트리밍 · MediaRecorder 인앱 녹음"),
    (ORANGE,   "패키징",      "Capacitor 6 + Gradle",
     "Android APK / AAB · 아이콘 · 스플래시 · 상태바 · 마이크 권한"),
    (CYAN,     "배포",        "Render · GitHub 자동 배포",
     "Render Free (15분 슬립) · Play Store 검토 중 · APK 사이드로드"),
]

col1, col2 = 0.4, 6.85
for i, (c, layer, tech, desc) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = col1 if col == 0 else col2
    y = 2.28 + row * 1.1
    rect(s, x, y, 6.1, 1.0, fill=NAVY_CARD)
    rect(s, x, y, 0.1, 1.0, fill=c)
    tb(s, x + 0.2, y + 0.07, 1.8, 0.38, layer, size=12, color=c, bold=True)
    tb(s, x + 2.1, y + 0.07, 3.85, 0.38, tech, size=13, color=WHITE, bold=True)
    tb(s, x + 0.2, y + 0.55, 5.8, 0.38, desc, size=11, color=GRAY)

footer_navy_bar(s, "3차 백엔드 → 4차 프론트엔드 → 5차 인증·배포 → 6차 패키징 완성")


# ════════════════════════════════════════════════════════════
# SLIDE 11 — 런칭 로드맵 & 수익 구조
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)
slide_title_navy(s, "런칭 로드맵 & 수익 구조",
                 "Phase 1 학기 내 완성(완료) → Phase 2 배포 → Phase 3 B2B 확장")

# Phase 카드 3개
phases = [
    (GREEN,    "Phase 1",  "학기 내 앱 완성",  "완료",
     ["✓ 백엔드 MVP (3차)",
      "✓ 웹 프론트엔드 (4차)",
      "✓ 인증 · 배포 (5차)",
      "✓ APK 패키징 (6차)",
      "✓ Play Store 등록 (6차)",
      "✓ 베타 테스트 완료"]),
    (SKY_BLUE, "Phase 2",  "배포 & 포트폴리오",  "진행 중",
     ["Render 배포 LIVE 유지",
      "Play Store 승인 후 링크 공유",
      "GitHub README 정비",
      "아키텍처 다이어그램 추가",
      "포트폴리오 아카이브 정리",
      "—"]),
    (ORANGE,   "Phase 3",  "파이프라인 확장",  "장기",
     ["감사·실사: 인터뷰 녹음 → 리스크 보고서",
      "컨설팅: 고객 인터뷰 → 분석 초안",
      "기업 교육: 사내 강의 → 지식 베이스",
      "법무: 회의록·계약서 → 검토 요약",
      "핵심: STT + Context + LLM 구조 재사용",
      "—"]),
]

for i, (c, phase, title, status, items) in enumerate(phases):
    x = 0.4 + i * 4.3
    rect(s, x, 1.5, 4.1, 4.05, fill=NAVY_CARD)
    rect(s, x, 1.5, 4.1, 0.1, fill=c)
    badge(s, x + 2.85, 1.64, status, c)
    tb(s, x + 0.15, 1.64, 2.7, 0.32, phase, size=12, color=c, bold=True)
    tb(s, x + 0.15, 1.98, 3.9, 0.46, title, size=16, color=WHITE, bold=True)
    rect(s, x + 0.15, 2.5, 3.8, 0.04, fill=c)
    for j, item in enumerate(items):
        tb(s, x + 0.15, 2.6 + j * 0.47, 3.9, 0.42, item, size=11, color=LIGHT_GRAY)

# 수익 구조 하단 (3열)
tb(s, 0.4, 5.68, 12.5, 0.38, "수익 구조 — 구독 모델", size=15, color=SKY_BLUE, bold=True)

plans = [
    (LIGHT_GRAY, CHARCOAL, "Free",  "₩0 / 월",
     ["월 10회 노트 생성", "강의 등록 3개", "기본 STT 품질"]),
    (SKY_BLUE,   WHITE,    "Pro",   "₩9,900 / 월",
     ["무제한 노트 생성", "강의 등록 무제한", "고품질 STT", "우선 처리"]),
    (NAVY_CARD,  WHITE,    "Team",  "₩29,900 / 월",
     ["Pro 기능 전체", "팀원 5명", "노트 공유", "관리자 대시보드"]),
]
for i, (bg, tc, name, price, feats) in enumerate(plans):
    px = 0.4 + i * 4.3
    rect(s, px, 6.1, 4.1, 1.32, fill=bg)
    tb(s, px + 0.15, 6.13, 2.0, 0.34, name, size=14, color=tc, bold=True)
    tb(s, px + 2.1, 6.13, 1.85, 0.34, price, size=12, color=tc)
    feat_txt = "  ·  ".join(feats)
    tb(s, px + 0.15, 6.48, 3.85, 0.88, feat_txt, size=10, color=tc)

footer_sky(s, "B2C 대학생 구독 → B2B 기관 라이선스 → 파이프라인 SaaS  |  경영전략 팀 프로젝트 9팀")


# ════════════════════════════════════════════════════════════
# SLIDE 12 — 마무리 + Q&A
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
navy_bg(s)

rect(s, 0.85, 1.3, 0.07, 5.0, fill=SKY_BLUE)

tb(s, 1.05, 1.3, 11, 0.55, "감사합니다.", size=42, color=WHITE, bold=True)
tb(s, 1.05, 1.95, 11, 0.45, "Q & A", size=28, color=SKY_BLUE, bold=True)

rect(s, 1.05, 2.55, 8.5, 0.05, fill=NAVY_CARD)

# 정리 3줄
summary = [
    "렉처노트는 음성·PDF·필기를 AI가 통합해 구조화된 강의 노트를 자동 생성하는 서비스입니다.",
    "1~6차 발표를 통해 아이디어 기획부터 Android APK 출시까지 13주를 완주했습니다.",
    "Phase 2 배포 완료 후 B2B 파이프라인 확장으로 에듀테크 시장을 공략하겠습니다.",
]
for i, line in enumerate(summary):
    tb(s, 1.05, 2.72 + i * 0.52, 11, 0.46, line, size=14, color=LIGHT_GRAY)

# LIVE DEMO 박스
rect(s, 1.05, 4.35, 7.5, 0.9, fill=NAVY_CARD)
tb(s, 1.2, 4.38, 2.2, 0.38, "LIVE DEMO", size=12, color=SKY_BLUE, bold=True)
tb(s, 1.2, 4.72, 7.2, 0.42,
   "lecture-note-2cb6.onrender.com",
   size=16, color=WHITE, bold=True)

# Play Store 박스
rect(s, 1.05, 5.38, 7.5, 0.9, fill=NAVY_CARD)
tb(s, 1.2, 5.41, 2.5, 0.38, "PLAY STORE", size=12, color=GREEN, bold=True)
tb(s, 1.2, 5.75, 7.2, 0.42,
   "com.kbslab.lecturenote  —  검토 예정",
   size=15, color=WHITE, bold=True)

# 팀 정보
tb(s, 1.05, 6.4, 8.5, 0.4,
   "이병헌 (팀장) · 김도환 · 최철원  |  경영전략 팀 프로젝트 9팀",
   size=14, color=GRAY)

# QR 영역 (우측)
rect(s, 10.0, 2.5, 3.0, 3.0, fill=WHITE, line=GRAY)
tb(s, 10.0, 2.5, 3.0, 3.0, "QR\n코드\n(배포 후\n삽입)",
   size=14, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 10.0, 5.6, 3.0, 0.38, "앱 다운로드 QR", size=11, color=GRAY, align=PP_ALIGN.CENTER)

footer_sky(s)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"\n[OK] 저장 완료: {OUTPUT_PATH}")
print(f"   슬라이드 수: {len(prs.slides)}장")
print(f"\n[사용법]")
print(f"   pip install python-pptx  (설치 안 된 경우)")
print(f"   python src/make_pptx_6th.py")
